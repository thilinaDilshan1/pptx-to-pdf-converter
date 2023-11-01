import os
from pptx import Presentation
from reportlab.pdfgen import canvas

def convert_single_pptx_to_pdf(pptx_path, pdf_path):
    prs = Presentation(pptx_path)
    pdf = canvas.Canvas(pdf_path)

    for slide_number, slide in enumerate(prs.slides):
        pdf.drawString(100, 800, f"Slide {slide_number + 1}")

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text
                pdf.drawString(100, 800 - (shape.top * 72), text)

        pdf.showPage()

    pdf.save()

def convert_bulk_pptx_to_pdf(folder_path, output_folder):
    pptx_files = [f for f in os.listdir(folder_path) if f.endswith(".pptx")]

    for pptx_file in pptx_files:
        pptx_path = os.path.join(folder_path, pptx_file)
        pdf_file = os.path.splitext(pptx_file)[0] + ".pdf"
        pdf_path = os.path.join(output_folder, pdf_file)

        convert_single_pptx_to_pdf(pptx_path, pdf_path)

# Example usage
conversion_type = input("Choose conversion type (1 for single, 2 for bulk): ")

if conversion_type == '1':
    pptx_path = input("Enter the path to the PowerPoint file: ")
    pdf_path = input("Enter the path for the PDF output: ")
    convert_single_pptx_to_pdf(pptx_path, pdf_path)

elif conversion_type == '2':
    folder_path = input("Enter the path to the folder containing PowerPoint files: ")
    output_folder = input("Enter the path for the output folder: ")
    convert_bulk_pptx_to_pdf(folder_path, output_folder)

else:
    print("Invalid choice. Please enter 1 or 2.")
